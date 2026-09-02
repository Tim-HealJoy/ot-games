#!/usr/bin/env python3
"""產生瀏覽器端流程測試用的 PPTX 與 PDF 樣本。

PPTX：兩張投影片，各放一張有浮水印的圖（一 PNG 一 JPEG）。
PDF ：兩頁，第 1 頁放有浮水印的圖、第 2 頁純文字（用來驗證文字層是否被保留）。
"""
import os
import shutil

HERE = os.path.dirname(os.path.abspath(__file__))
OUT = os.path.join(HERE, 'fixtures')
os.makedirs(OUT, exist_ok=True)

PNG_WM = os.path.join(HERE, 'samples', 'nano_96_wm.png')      # 星星浮水印
JPG_WM = os.path.expanduser('~/Downloads/Gemini_Generated_Image_yqmsxjyqmsxjyqms.jpeg')  # Notebook 文字

# ── PPTX ────────────────────────────────────────────────────────────────
from pptx import Presentation
from pptx.util import Inches, Pt

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

for path, title in [(PNG_WM, 'Nano Banana 星星'), (JPG_WM, 'Gemini Notebook 文字')]:
    slide = prs.slides.add_slide(blank)
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(8), Inches(0.7))
    box.text_frame.text = title
    box.text_frame.paragraphs[0].runs[0].font.size = Pt(28)
    slide.shapes.add_picture(path, Inches(0.5), Inches(1.1), width=Inches(9))

pptx_path = os.path.join(OUT, 'watermark_test.pptx')
prs.save(pptx_path)
print('PPTX →', pptx_path, os.path.getsize(pptx_path), 'bytes')

# 確認 media 內容
import zipfile
with zipfile.ZipFile(pptx_path) as z:
    media = [n for n in z.namelist() if n.startswith('ppt/media/')]
    print('  ppt/media/：', media)

# ── PDF ─────────────────────────────────────────────────────────────────
import fitz

# 第 1 頁：整頁滿版的浮水印圖。頁面尺寸（pt）＝圖片像素，也就是一般
# 「圖片轉 PDF」工具預設的 72dpi 換算——這才是使用者真正會拿來清的 PDF 形狀。
# 圖片若像原本那樣內縮留白，星星就不在「頁面」的右下角，逐頁 render 的做法本來就抓不到。
from PIL import Image as _PILImage
_w, _h = _PILImage.open(PNG_WM).size
doc = fitz.open()
page1 = doc.new_page(width=_w, height=_h)
page1.insert_image(fitz.Rect(0, 0, _w, _h), filename=PNG_WM)

page2 = doc.new_page(width=_w, height=_h)
page2.insert_text((60, 90), 'Page 2 - pure text, no image, no watermark.',
                  fontsize=20, fontname='helv')
page2.insert_text((60, 130), 'This page must keep its selectable text layer.',
                  fontsize=14, fontname='helv')
page2.insert_text((60, 160), 'If the tool rasterises it, the test fails.',
                  fontsize=14, fontname='helv')

pdf_path = os.path.join(OUT, 'watermark_test.pdf')
doc.save(pdf_path)
doc.close()
print('PDF  →', pdf_path, os.path.getsize(pdf_path), 'bytes')

# 順便把單張圖樣本複製一份到 fixtures，方便瀏覽器上傳
shutil.copy(JPG_WM, os.path.join(OUT, 'notebook_sample.jpeg'))
shutil.copy(PNG_WM, os.path.join(OUT, 'nano96_sample.png'))
print('圖片樣本已複製到 fixtures/')
