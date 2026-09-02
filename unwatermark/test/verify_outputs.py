#!/usr/bin/env python3
"""驗證瀏覽器實測產出的檔案（test/out/ 內）。

對應規格書測試 5（PPTX）與 6（PDF），另加圖片輸出的像素驗證。
"""
import os
import zipfile
import sys

HERE = os.path.dirname(os.path.abspath(__file__))
OUT = os.path.join(HERE, 'out')
FIX = os.path.join(HERE, 'fixtures')

ok_all = True


def head(t):
    print('\n' + t)


def check(cond, msg):
    global ok_all
    print(('  ✅ ' if cond else '  ❌ ') + msg)
    if not cond:
        ok_all = False
    return cond


# ── 5. PPTX ─────────────────────────────────────────────────────────────
head('## 測試 5：PPTX（瀏覽器輸出）')
pptx = os.path.join(OUT, 'watermark_test_已清除.pptx')
if not os.path.exists(pptx):
    check(False, '找不到 ' + pptx)
else:
    with zipfile.ZipFile(pptx) as z:
        bad = z.testzip()
        check(bad is None, 'ZIP 結構完整（testzip 無損毀項目）')
        names = z.namelist()
        check('[Content_Types].xml' in names, '[Content_Types].xml 存在')
        media = sorted(n for n in names if n.startswith('ppt/media/') and not n.endswith('/'))
        check(media == ['ppt/media/image1.png', 'ppt/media/image2.jpg'],
              '圖片檔名與副檔名未被更動：' + str(media))
    with zipfile.ZipFile(os.path.join(FIX, 'watermark_test.pptx')) as z0, \
         zipfile.ZipFile(pptx) as z1:
        n0, n1 = set(z0.namelist()), set(z1.namelist())
        check(n0 == n1, '檔案清單與原檔完全一致（沒有多／少檔）')
        xml_changed = [n for n in sorted(n0)
                       if n.endswith('.xml') or n.endswith('.rels')
                       if z0.read(n) != z1.read(n)]
        check(not xml_changed, 'XML／rels 全部逐位元組未更動' +
              ('，但發現：' + str(xml_changed) if xml_changed else ''))
        img_changed = [n for n in sorted(media) if z0.read(n) != z1.read(n)]
        check(len(img_changed) == 2, '兩張圖片內容都被替換過：' + str(img_changed))

    try:
        from pptx import Presentation
        prs = Presentation(pptx)
        check(len(prs.slides) == 2, 'python-pptx 可正常開啟，投影片數 = %d' % len(prs.slides))
        pics = sum(1 for s in prs.slides for sh in s.shapes if sh.shape_type == 13)
        check(pics == 2, '兩張圖片仍在版面上（shape_type=PICTURE 共 %d 個）' % pics)
    except Exception as e:
        check(False, 'python-pptx 讀取失敗：%s' % e)

# ── 6. PDF ──────────────────────────────────────────────────────────────
head('## 測試 6：PDF（瀏覽器輸出）')
pdf = os.path.join(OUT, 'watermark_test_已清除.pdf')
if not os.path.exists(pdf):
    check(False, '找不到 ' + pdf)
else:
    import fitz
    d0 = fitz.open(os.path.join(FIX, 'watermark_test.pdf'))
    d1 = fitz.open(pdf)
    check(d1.page_count == 2, '頁數 = %d（應為 2）' % d1.page_count)

    p1 = d1[0]
    t1 = p1.get_text().strip()
    imgs1 = p1.get_images(full=True)
    check(len(imgs1) == 1 and t1 == '',
          '第 1 頁（有浮水印的圖）已轉為單張圖片：圖 %d 張、文字 %d 字'
          % (len(imgs1), len(t1)))

    p2 = d1[1]
    t2 = p2.get_text().strip()
    t2src = d0[1].get_text().strip()
    check(len(t2) > 0, '第 2 頁仍有文字層，可選取（%d 字）' % len(t2))
    check(t2 == t2src, '第 2 頁文字與原檔逐字一致（未被轉成圖片）')
    check(len(p2.get_images(full=True)) == 0, '第 2 頁沒有被塞進任何圖片')

    # 第 1 頁的浮水印是否真的清掉：把兩份 PDF 的第 1 頁都 render 出來比對右下角
    import numpy as np
    def render(doc, page_no, zoom=1.0):
        pm = doc[page_no].get_pixmap(matrix=fitz.Matrix(zoom, zoom), alpha=False)
        return np.frombuffer(pm.samples, dtype=np.uint8).reshape(pm.height, pm.width, 3)
    a = render(d0, 0).astype(np.float32)
    b = render(d1, 0).astype(np.float32)
    if a.shape == b.shape:
        H, W, _ = a.shape
        # 星星落點：右下角 96px、邊距 64px（頁面 pt = 圖片 px，zoom=1）
        y0, y1 = H - 64 - 96, H - 64
        x0, x1 = W - 64 - 96, W - 64
        lum_a = a[y0:y1, x0:x1].mean(axis=2)
        lum_b = b[y0:y1, x0:x1].mean(axis=2)
        # 浮水印是加白，清掉後該區平均亮度應下降
        check(lum_b.mean() < lum_a.mean() - 1.0,
              '第 1 頁星星區平均亮度由 %.2f 降到 %.2f（白色浮水印已被移除）'
              % (lum_a.mean(), lum_b.mean()))
        # 星星以外的地方不該被動到太多
        outside_a = a[:H - 200, :W - 200]
        outside_b = b[:H - 200, :W - 200]
        mse = float(np.mean((outside_a - outside_b) ** 2))
        psnr = 99 if mse == 0 else 10 * np.log10(255.0 ** 2 / mse)
        check(psnr > 30, '浮水印以外區域 PSNR %.2f dB（>30 表示只有 render 重編碼的差異）' % psnr)
    else:
        check(False, '兩份 PDF 第 1 頁 render 尺寸不同：%s vs %s' % (a.shape, b.shape))
    d0.close(); d1.close()

# ── 圖片輸出 ────────────────────────────────────────────────────────────
head('## 附帶：瀏覽器圖片輸出的像素驗證')
import numpy as np
from PIL import Image
pairs = [
    ('browser_nano96_sample_已清除.png', 'nano_96_clean.png', 'samples'),
]
for out_name, ref_name, ref_dir in pairs:
    op = os.path.join(OUT, out_name)
    rp = os.path.join(HERE, ref_dir, ref_name)
    if not os.path.exists(op):
        check(False, '找不到 ' + op); continue
    a = np.array(Image.open(op).convert('RGB')).astype(np.float32)
    r = np.array(Image.open(rp).convert('RGB')).astype(np.float32)
    if a.shape != r.shape:
        check(False, '尺寸不符 %s vs %s' % (a.shape, r.shape)); continue
    # 星星區
    H, W, _ = a.shape
    y0, y1, x0, x1 = H - 64 - 96, H - 64, W - 64 - 96, W - 64
    err = np.abs(a[y0:y1, x0:x1] - r[y0:y1, x0:x1]).max()
    mse = float(np.mean((a[y0:y1, x0:x1] - r[y0:y1, x0:x1]) ** 2))
    psnr = 99 if mse == 0 else 10 * np.log10(255.0 ** 2 / mse)
    check(err <= 2, '%s：星星區與乾淨原圖最大誤差 %d（≤2）、PSNR %.2f dB'
          % (out_name, int(err), psnr))

zp = os.path.join(OUT, 'browser_all.zip')
if os.path.exists(zp):
    with zipfile.ZipFile(zp) as z:
        check(z.testzip() is None, '「全部下載（ZIP）」產出的壓縮檔結構完整')
        check(len(z.namelist()) == 4, 'ZIP 內含 %d 個檔案：%s'
              % (len(z.namelist()), z.namelist()))

print('\n總結：' + ('✅ 全部通過' if ok_all else '❌ 有項目未通過'))
sys.exit(0 if ok_all else 1)
